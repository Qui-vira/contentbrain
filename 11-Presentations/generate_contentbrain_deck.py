"""Generate ContentBrain System Overview presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(18, 18, 18)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(220, 38, 38)
GRAY = RGBColor(160, 160, 160)
DARK_GRAY = RGBColor(40, 40, 40)
CELL_BG = RGBColor(28, 28, 28)
HEADER_BG = RGBColor(50, 10, 10)


def dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK


def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE,
                 bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tf


def add_bullet_slide(prs, title_text, bullets, sub_bullets=None):
    """Add a slide with title and bullet points. sub_bullets is a dict of index -> list of sub-items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    dark_bg(slide)

    # Red accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.8),
                 title_text, size=32, bold=True, color=WHITE)

    # Thin line under title
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    if sub_bullets is None:
        sub_bullets = {}

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

        if i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.space_after = Pt(3)
                sr = sp.add_run()
                sr.text = f"       {sub}"
                sr.font.size = Pt(13)
                sr.font.color.rgb = GRAY
                sr.font.name = "Calibri"

    return slide


def add_table_slide(prs, title_text, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    dark_bg(slide)

    # Red accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.8),
                 title_text, size=32, bold=True, color=WHITE)

    # Line
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(9.0)
    tbl_height = Inches(0.4 * n_rows)
    top = Inches(1.4)
    left = Inches(0.5)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width, tbl_height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = RED
            p.font.bold = True
            p.font.name = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CELL_BG if i % 2 == 0 else DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"

    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE 1: Title ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    # Large red accent block
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(10), Inches(1.5),
                 "CONTENTBRAIN", size=54, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.0), Inches(2.8), Inches(10), Inches(1.0),
                 "AI-Powered Content-to-Revenue Engine", size=28, color=RED)
    add_text_box(slide, Inches(1.0), Inches(4.0), Inches(10), Inches(0.6),
                 "System Overview  |  @big_quiv  |  Quivira", size=18, color=GRAY)

    # Bottom line
    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(5.0), Inches(5), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(5.3), Inches(10), Inches(0.5),
                 "53 AI Skills  |  62 CLI Tools  |  72 Integration Guides  |  Full Automation Pipeline",
                 size=14, color=GRAY)

    # ==================== SLIDE 2: What Is ContentBrain ====================
    add_bullet_slide(prs, "What Is ContentBrain?", [
        "An Obsidian vault + Claude Code automation system powering the Quivira brand",
        "One-person content agency: research, create, produce, publish, trade, track",
        "53 slash commands covering the entire content + marketing + trading stack",
        "Niche: Web3, crypto, AI, trading",
        "Platforms: X, LinkedIn, TikTok, Instagram, Telegram, YouTube",
        "Voice: Smart, bold, streetwise, alpha-coded",
    ])

    # ==================== SLIDE 3: Vault Structure ====================
    add_table_slide(prs, "Vault Structure", ["Folder", "Purpose"], [
        ["00-Inbox", "Master strategy, KPIs, raw ideas"],
        ["01-Competitors", "Scraped competitor data & insights"],
        ["02-Hooks", "Hook libraries by goal, format, platform, topic"],
        ["03-Trends", "Trend scouting reports"],
        ["04-Patterns", "Content patterns extracted from competitors"],
        ["05-Frameworks", "Reusable content / sales / carousel frameworks"],
        ["06-Drafts", "Ready-to-post content: tweets, LinkedIn, TikTok, Telegram, PDFs"],
        ["07-Analytics", "Posting logs, hook performance, signal performance"],
        ["08-Templates + Media", "Voice/character/prompt libraries, video formats, media assets"],
        ["09-Skills", "Original skill definitions"],
        ["10-Niche-Knowledge", "Deep research: personal brand, graphic design, product specs"],
        ["11-Presentations", "Slide decks and presentation files"],
        ["skills/marketing-external", "26 imported marketing skills + tools library"],
    ], col_widths=[3.0, 6.0])

    # ==================== SLIDE 4: Original Skills (27) ====================
    add_bullet_slide(prs, "Original Skills (27 Commands)", [
        "Content Creation",
        "Business & Sales",
        "Publishing Pipeline",
        "Trading & Signals",
        "Research & Analytics",
    ], sub_bullets={
        0: ["/ghostwriter  /concept  /graphic-designer  /video-editor  /content-strategist"],
        1: ["/funnel-builder  /sales-closer  /community-manager  /operations-lead"],
        2: ["/post  /publish  /publish-update  /publish-status  /manychat-dm"],
        3: ["/scan-all  /scan-pair  /scan-custom  /forex-scan  /forex-pair  /forex-custom",
            "/technical-analyst  /send-signals  /send-signals-direct  /signal-tracker  /market-report"],
        4: ["/scrape-instagram  /data-analyst"],
    })

    # ==================== SLIDE 5: Imported Skills (26) ====================
    add_bullet_slide(prs, "Imported Marketing Skills (26 Commands)", [
        "CRO Suite (7)",
        "SEO Suite (5)",
        "Ads & Growth (5)",
        "Strategy & Ops (6)",
        "Content & Copy (3)",
    ], sub_bullets={
        0: ["/ab-test-setup  /form-cro  /onboarding-cro  /page-cro  /paywall-upgrade-cro  /popup-cro  /signup-flow-cro"],
        1: ["/ai-seo  /seo-audit  /schema-markup  /programmatic-seo  /site-architecture"],
        2: ["/ad-creative  /paid-ads  /cold-email  /referral-program  /free-tool-strategy"],
        3: ["/marketing-ideas  /marketing-psychology  /launch-strategy  /pricing-strategy  /product-marketing-context  /revops"],
        4: ["/copy-editing  /competitor-alternatives  /analytics-tracking"],
    })

    # ==================== SLIDE 6: Merge Summary ====================
    add_table_slide(prs, "Merge Summary", ["Metric", "Count"], [
        ["External skills evaluated", "33"],
        ["Skills added", "26"],
        ["Skills skipped (overlap)", "7"],
        ["Existing files modified", "0"],
        ["Total commands after merge", "53"],
    ], col_widths=[5.0, 4.0])

    # ==================== SLIDE 7: Skipped Skills ====================
    add_table_slide(prs, "Skipped Skills (Direct Overlap)", ["External Skill", "Covered By", "Reason"], [
        ["content-strategy", "content-strategist", "Same core function"],
        ["social-content", "ghostwriter + content-strategist", "Social posts + scheduling overlap"],
        ["copywriting", "ghostwriter + funnel-builder", "Marketing copy + landing pages"],
        ["email-sequence", "funnel-builder", "Email sequences + drip campaigns"],
        ["lead-magnets", "funnel-builder", "Lead magnet creation"],
        ["churn-prevention", "community-manager", "Churn prevention sequences"],
        ["sales-enablement", "sales-closer", "Objection handling + pitches"],
    ], col_widths=[2.5, 3.5, 3.0])

    # ==================== SLIDE 8: Tools Library ====================
    add_bullet_slide(prs, "Marketing Tools Library", [
        "62 CLI Scripts (.js)",
        "72 Integration Guides (.md)",
        "Composio Marketing Integrations",
        "Key platforms covered:",
    ], sub_bullets={
        0: ["Wrapper scripts for direct API interaction with marketing platforms"],
        1: ["Setup and usage documentation for each platform"],
        2: ["Marketing tool orchestration layer"],
        3: ["GA4, Ahrefs, Semrush, Meta Ads, Google Ads, Mailchimp, HubSpot, Klaviyo",
            "TikTok Ads, LinkedIn Ads, Buffer, Segment, Mixpanel, Amplitude",
            "Stripe, Shopify, Intercom, Customer.io, Postmark, SendGrid, Zapier"],
    })

    # ==================== SLIDE 9: Automation Scripts ====================
    add_table_slide(prs, "Automation Scripts", ["Script", "Function"], [
        ["binance_ta_runner.py", "Crypto technical analysis scanner"],
        ["forex_ta_runner.py", "Forex technical analysis scanner"],
        ["unified_auto_scanner.py", "Automated signal pipeline"],
        ["signal_monitor.py", "Signal monitoring + alerts"],
        ["polymarket_bot.py", "Polymarket event tracking"],
        ["polymarket_scanner.py", "Polymarket opportunity scanner"],
        ["market_data.py", "Market data feeds"],
        ["fetch_alerts.py", "TradingView alert fetcher"],
        ["post_reel.py", "Instagram reel posting"],
        ["deploy-manychat-v2.py", "ManyChat automation deployment"],
    ], col_widths=[3.5, 5.5])

    # ==================== SLIDE 10: Infrastructure ====================
    add_table_slide(prs, "Infrastructure", ["Service", "Role"], [
        ["Supabase", "Database — signal dedup, tracking, state"],
        ["Railway", "Single deployment (quivira-signal-bot, BOT_MODE=full)"],
        ["Telegram", "Signal distribution + community management"],
        ["Notion", "Content calendar + publish pipeline"],
        ["Typefully", "X + LinkedIn posting"],
        ["fal.ai", "AI image generation (Nano Banana)"],
        ["MiniMax", "Voice cloning (primary for @big_quiv)"],
        ["ElevenLabs", "Foreign character voices only"],
        ["Apify", "Social media scraping (all platforms)"],
        ["ManyChat", "Instagram DM automation"],
    ], col_widths=[3.0, 6.0])

    # ==================== SLIDE 11: Data Layer ====================
    add_bullet_slide(prs, "Data Layer", [
        "Pre-scraped raw data (JSON) across platforms and niches:",
        "Instagram: competitors, crypto/AI reels, ghostwriting, visual hooks",
        "TikTok: ghostwriting, graphic design, visual hooks",
        "YouTube: ghostwriting, graphic design, visual hooks, transcripts",
        "Specialized data: copywriting intake, sales funnel intake, video production intake",
        "Trading data: binance_ta_summary.json, forex_ta_summary.json, polymarket_signals_temp.json",
    ])

    # ==================== SLIDE 12: Content Strategy ====================
    add_bullet_slide(prs, "Content Strategy Engine", [
        "Hook Formula: Callout-Flex-Reveal",
        "Content Mix: 30% value, 25% story, 20% community, 15% promo, 10% hot takes",
        "Every post serves one goal: sales, reach, leads, authority, or community",
        "Dream Client Formula:",
        "Movement Multipliers:",
        "Batch create on weekends, post throughout the week",
    ], sub_bullets={
        3: ["Value (teach something worth saving) + Relatability (speak their language)",
            "Clarity (no fluff, straight signal) + Positioning (reinforce authority)"],
        4: ["Polarizing stands, shared mythology, common enemy, energetic delivery"],
    })

    # ==================== SLIDE 13: Closing ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(2.0), Inches(10), Inches(1.5),
                 "CONTENTBRAIN", size=54, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.0), Inches(3.3), Inches(10), Inches(1.0),
                 "53 Commands. Full Stack. One System.", size=28, color=RED)
    add_text_box(slide, Inches(1.0), Inches(4.5), Inches(10), Inches(0.6),
                 "Research  >  Strategy  >  Create  >  Produce  >  Publish  >  Track  >  Trade",
                 size=18, color=GRAY)

    shape = slide.shapes.add_shape(1, Inches(1.0), Inches(5.5), Inches(5), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(5.8), Inches(10), Inches(0.5),
                 "@big_quiv  |  Quivira  |  Built with Claude Code",
                 size=14, color=GRAY)

    # Save
    out = "C:/Users/Bigquiv/onedrive/desktop/contentbrain/11-Presentations/ContentBrain-System-Overview.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
