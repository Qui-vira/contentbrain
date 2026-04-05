"""Generate Outlier AI PPTX Presentation."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT_DIR = os.path.join(os.path.dirname(__file__), '..', '06-Drafts')
OUT_PATH = os.path.join(OUT_DIR, 'outlier-ai-presentation.pptx')

DARK = RGBColor(0x1a, 0x1a, 0x2e)
RED = RGBColor(0xe9, 0x45, 0x60)
GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title_text, bullets=None, subtitle=None, is_title_slide=False):
    if is_title_slide:
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # Subtitle
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.333), Inches(1.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0xaa, 0xaa, 0xaa)
            p2.alignment = PP_ALIGN.CENTER
        return slide

    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Red accent bar at top
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True

    # Bullets
    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = b
            p2.font.size = Pt(18)
            p2.font.color.rgb = GRAY
            p2.space_after = Pt(10)
            p2.level = 0

    return slide

def add_table_slide(title_text, headers, rows):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Red bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.5),
        Inches(11.7), Inches(5)
    ).table

    # Style header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = GRAY

    return slide

# BUILD SLIDES
# Slide 1: Title
add_slide(
    "Outlier AI",
    is_title_slide=True,
    subtitle="Complete Registration & Earning Guide (2026)\n\n13 Twitter sources + 8 web sources | April 2026"
)

# Slide 2: What is Outlier?
add_slide("What is Outlier?", [
    "• AI training platform operated by Scale AI ($14B valuation)",
    "• Pays freelancers to evaluate, correct, and improve AI responses",
    "• $100M+ paid to 40,000+ contributors globally",
    "• 40+ domains: writing, coding, STEM, languages, law, medicine",
    "• Sister platform to Remotasks (both owned by Scale AI)",
])

# Slide 3: What tasks do you do?
add_table_slide("What Tasks Do You Do?",
    ['Task Type', 'What You Do', 'Example'],
    [
        ['Response Evaluation', 'Compare AI answers, pick the better one', '"Which response is more accurate?"'],
        ['Content Rewriting', 'Fix AI text for accuracy and clarity', 'Correct factual errors'],
        ['Coding Tasks', 'Review AI code, debug, improve', 'Fix a Python function'],
        ['Fact-Checking', 'Verify AI claims against real data', 'Check statistics'],
        ['Prompt Engineering', 'Write questions to test AI limits', 'Create tricky math problems'],
        ['Data Annotation', 'Label and categorize content', 'Tag helpful vs harmful'],
    ]
)

# Slide 4: Requirements
add_slide("Requirements", [
    "MUST HAVE:",
    "• Age: 18+",
    "• Education: Associate degree minimum (bachelor's+ preferred)",
    "• Valid government ID (passport or driver's license)",
    "• Updated resume + LinkedIn profile",
    "• Laptop/desktop with stable internet",
    "",
    "GETS YOU HIGHER PAY:",
    "• Graduate degree (Master's, PhD)",
    "• STEM, coding, legal, or medical expertise",
    "• GitHub portfolio or published academic work",
])

# Slides 5-11: Step by step
steps = [
    ("Step 1: Go to outlier.ai", [
        "• Visit https://app.outlier.ai/en/expert/signup",
        "• Click 'Create an Account'",
        "• Sign up with email or Google account",
        "",
        "TIP: 'Creating an account is one thing —",
        "passing onboarding is a whole different game.' — @Fra_nkXBT",
    ]),
    ("Step 2: Create Account & Profile", [
        "• Enter: full name, phone, DOB, country, address",
        "• Upload your current resume (auto-extracts skills)",
        "• Connect your LinkedIn profile (cross-referenced)",
        "",
        "NIGERIA: Platform rejects +234 numbers",
        "• Use paid VPN + US virtual number",
        "• OR use RDP (preferred — real machine fingerprint)",
        "• Recommended proxies: NSocks, CleanLTE",
    ]),
    ("Step 3: Identity Verification", [
        "• Uses Persona (same service banks use)",
        "• Take clear photos of ID front and back",
        "• Take a live selfie (well-lit, face centered)",
        "• Processing: 1-2 business days",
        "",
        "COMMON FAILURES:",
        "• Expired ID",
        "• Blurry photos",
        "• ID details don't match profile info",
    ]),
    ("Step 4: Select Skills", [
        "• Choose maximum 10 skills",
        "• Only pick skills you can genuinely demonstrate",
        "• You WILL be tested on them",
        "",
        "HIGH DEMAND IN 2026:",
        "• Advanced mathematics, Python, JavaScript",
        "• Legal analysis, medical reasoning",
        "• Data science, linguistics, technical writing",
    ]),
    ("Step 5: Take Qualification Test", [
        "• General Reasoning Screening",
        "• Must score 80% or higher",
        "• Tests analytical thinking, reading comprehension, logic",
        "• Often ONE attempt only — prepare carefully",
        "",
        "APRIL 2026 UPDATE:",
        "• Aether assessment reportedly temporarily suspended",
        "• Check back regularly if unavailable",
    ]),
    ("Step 6: Complete Onboarding", [
        "• Project-specific training (may be unpaid)",
        "• Read ALL guidelines thoroughly",
        "• Watch training videos completely",
        "",
        "TIMELINE:",
        "• Active onboarding: 30-90 minutes",
        "• ID verification: 1-2 days",
        "• First project assignment: days to weeks",
    ]),
    ("Step 7: Start Earning", [
        "• Browse available projects on your dashboard",
        "• Each listing shows: task, time, pay rate",
        "• Accept tasks matching your verified skills",
        "• Work daily cap: 6, 8, or 10 hours (varies by project)",
        "• Daily reset: ~4am",
        "",
        "'Accuracy matters more than speed.",
        "Mess up too much and they'll kick you out.' — @iamveektoria_",
    ]),
]

for title, bullets in steps:
    add_slide(title, bullets)

# Slide 12: Pay rates table
add_table_slide("Pay Rates",
    ['Category', 'Hourly Rate'],
    [
        ['Entry-level tasks', '$12-$20/hr'],
        ['Writing/content evaluation', '$15-$35/hr'],
        ['Math/STEM review', '$20-$50/hr'],
        ['Coding/technical', '$25-$50/hr'],
        ['Software engineering', '~$43/hr'],
        ['Legal review', '$50-$100/hr'],
        ['Medical review', '$50-$100/hr'],
        ['US average (all contributors)', '~$31/hr (Glassdoor)'],
    ]
)

# Slide 13: How you get paid
add_slide("How You Get Paid", [
    "PAYMENT METHODS:",
    "• PayPal (most common, works in Nigeria)",
    "• Airtm (good for African users)",
    "• ACH Bank Transfer (US only)",
    "",
    "SCHEDULE:",
    "• Processed every Tuesday",
    "• For work done previous Tuesday-Monday (midnight UTC)",
    "• Minimum payout: $10",
    "• Funds arrive: 1-3 business days",
    "",
    "NIGERIA TIP: Use domiciliary (USD) account for PayPal to avoid conversion losses",
])

# Slide 14: Tips
add_slide("Tips to Get Accepted", [
    '1. Use RDP over VPN (real US fingerprint) — @ObaDeleke',
    '2. Accuracy > speed always — quality score determines everything',
    '3. Use solid proxy from Nigeria: NSocks, CleanLTE',
    '4. Keep resume current — system auto-extracts skills',
    '5. Connect strong LinkedIn — cross-referenced for credibility',
    '6. Only select 5-10 genuine skills — fakes get caught',
    '7. Highlight STEM/law/medicine — higher pay + priority',
    '8. Score 80%+ on reasoning test — non-negotiable',
    '9. Read ALL guidelines before any assessment',
    '10. Be patient — tasks can take weeks to arrive',
])

# Slide 15: Common mistakes
add_slide("Common Mistakes", [
    "REGISTRATION:",
    "• Claiming skills you can't demonstrate → locked out",
    "• Using expired or blurry ID → verification fails",
    "• Rushing reasoning test → score below 80%, no retake",
    "",
    "WHILE WORKING:",
    "• Speed over accuracy → quality drops, account banned",
    "• Using ChatGPT/AI tools → detected and banned immediately",
    "• Ignoring project guidelines → flagged and restricted",
    "• Letting large balance build up → risk if account deactivated",
    "",
    "PLATFORM RISKS: Silent removal, no appeals process, automated reviews",
])

# Slide 16: Realistic earnings
add_table_slide("Realistic Earnings",
    ['Scenario', 'Hours/Week', 'Rate', 'Monthly'],
    [
        ['Casual', '5-10 hrs', '$15/hr', '$300-$600'],
        ['Consistent', '15-20 hrs', '$20/hr', '$1,200-$1,600'],
        ['Active', '30-40 hrs', '$25/hr', '$3,000-$4,000'],
        ['Specialist', '20-30 hrs', '$40-60/hr', '$3,200-$7,200'],
    ]
)

# Slide 17: Get started
add_slide(
    "Get Started",
    is_title_slide=True,
    subtitle=(
        "Sign up: https://app.outlier.ai/en/expert/signup\n\n"
        "Remember: Accuracy > Speed | Prepare before applying | Combine with other platforms\n\n"
        "Sources: 13 X/Twitter posts + 8 web sources | April 2026"
    )
)

prs.save(OUT_PATH)
print(f"PPTX saved to: {OUT_PATH}")
